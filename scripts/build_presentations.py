#!/usr/bin/env python3
"""
Silver Medal — presentation generator (v2026).

Builds four modern, branded 16:9 decks from the language in
silvermedal.net/server/assets/presentations/OUTLINES_v2026.md:

  * SilverMedal_Coding_v2026.pptx      (Red Rock Code)
  * SilverMedal_Writing_v2026.pptx     (Writers Guild)
  * SilverMedal_Healthcare_v2026.pptx  (Life Science Balance)
  * SilverMedal_Overview_v2026.pptx    (Silver Medal ecosystem)

New files are written alongside the existing decks; nothing is overwritten.

Run:
    python3 -m venv .venv-pptx
    .venv-pptx/bin/pip install python-pptx pillow
    .venv-pptx/bin/python scripts/build_presentations.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Paths & palette
# ---------------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, ".."))
PUB = os.path.join(ROOT, "silvermedal.net", "server", "public")
IMG = os.path.join(PUB, "assets", "images")
LSB = os.path.join(PUB, "images", "lsb")
# Write into the PUBLIC tree so the decks are web-accessible at
# /assets/presentations/... (served by a-silvermedal-static / assets-fallback).
OUT = os.path.join(PUB, "assets", "presentations")

INK = RGBColor(0x0F, 0x1F, 0x33)
BLUE = RGBColor(0x25, 0x63, 0xEB)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
TEAL = RGBColor(0x0D, 0x94, 0x88)
AMBER = RGBColor(0xD9, 0x77, 0x06)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MIST = RGBColor(0xCB, 0xD5, 0xE1)
SLATE = RGBColor(0x33, 0x41, 0x55)
CLOUD = RGBColor(0xF5, 0xF7, 0xFB)

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

FONT = "Calibri"


def img(*parts):
    """Return an image path if it exists, else None (slides skip missing art)."""
    p = os.path.join(*parts)
    return p if os.path.exists(p) else None


# Curated, verified-on-disk imagery per deck.
LOGO = img(IMG, "silvermedal", "silver_medal_emblem_transparent_bg.png")
FOUR_ROBES = img(IMG, "silvermedal", "silvermedal_four_robes_0.png")


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bg(slide, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    _fill(r, color)
    r.shadow.inherit = False
    # push to back
    sp = r._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return r


def band(slide, x, y, w, h, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _fill(r, color)
    r.shadow.inherit = False
    return r


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def setpar(p, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
           space_after=6, font=FONT):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return p


def picture_cover(slide, path, x, y, w, h):
    """Add a picture scaled to fill (cover) the box, cropped to the box."""
    if not path:
        return None
    try:
        from PIL import Image
        iw, ih = Image.open(path).size
    except Exception:
        pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
        return pic
    box_ratio = w / h
    img_ratio = iw / ih
    pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
    if img_ratio > box_ratio:
        # image wider -> crop left/right
        crop = (1 - box_ratio / img_ratio) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        crop = (1 - img_ratio / box_ratio) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    return pic


# ---------------------------------------------------------------------------
# Slide templates
# ---------------------------------------------------------------------------
def title_slide(prs, title, kicker, subtitle, accent, image=None, logo=LOGO):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, INK)
    # accent side panel with image
    panel_w = Inches(4.6)
    band(slide, EMU_W - panel_w, 0, panel_w, EMU_H, accent)
    if image:
        picture_cover(slide, image, EMU_W - panel_w, 0, panel_w, EMU_H)
        # translucent accent veil for legibility
        veil = band(slide, EMU_W - panel_w, 0, panel_w, EMU_H, accent)
        veil.fill.fore_color.rgb = accent
        veil.fill.transparency = 0  # solid fallback; kept subtle by size
        veil.line.fill.background()
        # make the veil a thin bottom gradient-like strip instead of covering all
        veil.height = Inches(2.2)
        veil.top = EMU_H - Inches(2.2)
    # accent rule
    band(slide, Inches(0.9), Inches(2.5), Inches(1.1), Inches(0.11), accent)
    if logo:
        slide.shapes.add_picture(logo, Inches(0.85), Inches(0.7),
                                 height=Inches(1.15))
    tb, tf = textbox(slide, Inches(0.9), Inches(2.8), Inches(7.4), Inches(3.4))
    setpar(tf.paragraphs[0], kicker, 16, accent, bold=True)
    setpar(tf.add_paragraph(), title, 46, WHITE, bold=True, space_after=10)
    setpar(tf.add_paragraph(), subtitle, 20, MIST, space_after=0)
    return slide


def content_slide(prs, kicker, title, bullets, accent, image=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, CLOUD)
    # header band
    band(slide, 0, 0, EMU_W, Inches(1.5), INK)
    band(slide, 0, Inches(1.5), EMU_W, Inches(0.08), accent)
    tb, tf = textbox(slide, Inches(0.9), Inches(0.28), Inches(11.5), Inches(1.1))
    setpar(tf.paragraphs[0], kicker, 13, accent, bold=True, space_after=2)
    setpar(tf.add_paragraph(), title, 30, WHITE, bold=True, space_after=0)

    has_img = bool(image)
    text_w = Inches(7.4) if has_img else Inches(11.5)
    tb, tf = textbox(slide, Inches(0.9), Inches(1.95), text_w, Inches(5.1),
                     anchor=MSO_ANCHOR.TOP)
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if isinstance(b, tuple):
            lead, rest = b
            setpar(p, lead, 18, INK, bold=True, space_after=3)
            r = p.add_run()
            r.text = "  " + rest
            r.font.size = Pt(18)
            r.font.color.rgb = SLATE
            r.font.name = FONT
        else:
            setpar(p, "•  " + b, 18, SLATE, space_after=8)
    if has_img:
        picture_cover(slide, image, Inches(8.55), Inches(2.05),
                      Inches(3.9), Inches(4.7))
    return slide


def closing_slide(prs, big, url, accent, image=None, logo=LOGO):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, INK)
    if image:
        picture_cover(slide, image, 0, 0, EMU_W, EMU_H)
        veil = band(slide, 0, 0, EMU_W, EMU_H, INK)
        veil.fill.fore_color.rgb = INK
        veil.fill.transparency = 0
        # emulate darken by shrinking to a gradient-ish overlay: keep full but
        # rely on ink; put text on top
    band(slide, Inches(5.9), Inches(3.05), Inches(1.5), Inches(0.12), accent)
    if logo:
        slide.shapes.add_picture(logo, Inches(6.06), Inches(1.5),
                                 height=Inches(1.3))
    tb, tf = textbox(slide, Inches(1.0), Inches(3.4), Inches(11.3), Inches(2.2),
                     anchor=MSO_ANCHOR.TOP)
    setpar(tf.paragraphs[0], big, 40, WHITE, bold=True, align=PP_ALIGN.CENTER,
           space_after=8)
    setpar(tf.add_paragraph(), url, 22, accent, bold=True,
           align=PP_ALIGN.CENTER, space_after=0)
    return slide


def new_prs():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    return prs


# ---------------------------------------------------------------------------
# Deck builders
# ---------------------------------------------------------------------------
def build_coding():
    prs = new_prs()
    ci = lambda n: img(IMG, "redrock", n)
    ei = lambda n: img(IMG, "extra", n)
    title_slide(
        prs, "Red Rock Code", "CODING · REASONING · DEVELOPING",
        "Project-based app development, taught from the ground up.  ·  redrockcode.com",
        RED, image=ci("rrcc_shots_0.png"))
    content_slide(prs, "WHAT MAKES US DIFFERENT", "A Standard-First Philosophy", [
        ("Standard ES6 / PWA / SPA", "module programming — from the ground up"),
        ("MVC & OOP", "principles woven through every project"),
        "<script type=\"module\"> + clean browser APIs — no framework mess",
        "Write the code once → ship to web, iOS, and Android as a PWA",
        "A bet we made in 2017 — and the right one",
    ], RED, image=ei("redrockcode_advanced_subjects_server_side.png"))
    content_slide(prs, "THE WORKBENCH", "Signature Projects", [
        ("Virtual Life", "every step of the animated life simulation"),
        ("User Management Page", "every step, MVC from scratch"),
        ("Bluesky", "a complete real-estate SPA (live: redrockbluesky.org)"),
        ("Responsive Web Design", "Epilogue, Industrious, Visualize, Interphase"),
        ("Classic capstones", "Platform Game, Personality Quiz, Quizzap Timer"),
    ], RED, image=ei("redrockcode_bluesky_project_responsive_web_capstone.png"))
    content_slide(prs, "THE CURRICULUM", "From 1000 to 3090", [
        ("Courses 1000 → 3090", "a certification issued for each completed term"),
        ("Stanford CS 106A+", "methodology, adapted and extended"),
        "Custom projects that go well beyond the FCC challenges",
        "Browse every implementation at redrockcode.com/educationMaterials",
    ], RED, image=ei("redrockcode_haverbeke_platform_game_project.png"))
    content_slide(prs, "ON THE HORIZON", "Robotics", [
        "CS 106A + 108 materials adapting toward a Robotics course",
        ("Micro-electronics", "Arduino, Raspberry Pi"),
        "Karel the Robot — meets real, physical moving parts",
        "Accreditation being pursued in the State of Utah",
    ], RED, image=ei("midjourney_RRCA_logo_1.png"))
    content_slide(prs, "WHY IT MATTERS", "Portable, Durable, Real", [
        "Portable PWAs to mobile — no React Native required",
        "HTML & CSS remain vital — as relevant now as in 2017",
        "We do teach React in the Code Camp (e.g. the Quizzap Quiz Timer)",
        ("Our mission", "build creators"),
    ], RED, image=ci("rrcc_shots_2.png"))
    closing_slide(prs, "Build. Solve. Elevate.", "redrockcode.com", RED)
    save(prs, "SilverMedal_Coding_v2026.pptx")


def build_writing():
    prs = new_prs()
    ei = lambda n: img(IMG, "extra", n)
    title_slide(
        prs, "Writers Guild", "WRITING · PUBLISHING · COLLABORATING",
        "writersguild.online  ·  writersguild.shop",
        PURPLE, image=ei("writers_guild_firefly_set2.png"))
    content_slide(prs, "A HOME FOR WRITERS", "Create, Publish, Keep Editing", [
        "Collaborate, publish, sell, and keep editing — multiple editions & formats",
        "Every author: their own site, their own brand — even their own domain",
        "Modular courseware and e-commerce, by default",
        "Built to meet the needs of any active contributor and co-creator",
    ], PURPLE, image=ei("writers_guild_firefly_set24_transparent.png"))
    content_slide(prs, "CLASSICS, REBORN", "Editions Side by Side", [
        "New editions of classics, compared side-by-side with annotations",
        ("Morals & Dogma — Acacia Concordance", "two transcriptions of Pike · 33 modules"),
        "Next: works with multiple historical editions worthy of modern review",
    ], PURPLE, image=ei("magus_midjourney.png"))
    content_slide(prs, "NEW WORKS", "The Seed Trilogy", [
        ("The 144,000", "beginning with \u201cThe Watchful One\u201d — the root of the Guild"),
        ("THE WORD", "a live-action trilogy (currently editing) · 45 modules"),
        "Written 100% by hand, then typed by human hands",
    ], PURPLE, image=ei("high-priestess_midjourney.png"))
    content_slide(prs, "OPENING THE DOOR TO AI", "Human & AI, Connected", [
        "AI invited as a co-author — spin-offs and \u201cforks\u201d that add new dimensions",
        "Keeping diverse forks connected, whether human- or AI-written",
        ("The Silver-Age Renaissance Series", "SARS"),
    ], PURPLE, image=ei("144000_animated_00001.gif"))
    content_slide(prs, "FOR EVERY CREATOR", "A Stage for All Voices", [
        "Screenplays and new theatrical productions — artists, actors, writers at work",
        "Professors publishing and selling their own textbooks",
        "Local clubs on the Quizzap core: courses, directories, even cookbooks",
    ], PURPLE, image=ei("writers_guild_firefly_set33.png"))
    content_slide(prs, "OUR MODELS OF EXCELLENCE", "Patience, Craft, and Joy", [
        "Stan Lee, George Lucas, Walt Disney",
        "Great universes take time to reach stardom",
        "We aim to emulate them in spirit and in deed",
    ], PURPLE, image=ei("writers_guild_firefly_set45.png"))
    closing_slide(prs, "Create. Share. Publish.", "writersguild.online", PURPLE)
    save(prs, "SilverMedal_Writing_v2026.pptx")


def build_healthcare():
    prs = new_prs()
    ei = lambda n: img(IMG, "extra", n)
    si = lambda n: img(IMG, "silvermedal", n)
    li = lambda n: img(LSB, n)
    title_slide(
        prs, "Life Science Balance", "ARTWORK · HEALING · WELLNESS",
        "lifesciencebalance.co",
        GREEN, image=ei("mandala_luminescent_0.png"))
    content_slide(prs, "A HOLISTIC MODEL", "Whole-Person Health", [
        ("Three-part holistic health certification", "free"),
        "High-class services and shops, in conjunction with partners",
        ("NCLEX-RN Test Prep 2026", "for aspiring nurses"),
    ], GREEN, image=li("lsb_tight_14.png"))
    content_slide(prs, "ART AS MEDICINE", "Creativity & Healing", [
        ("A growing Graphic Arts course", "new for 2026"),
        "Creating art — like writing — is fundamental to healing and balance",
        "A heavy, intentional artistic influence across the whole platform",
    ], GREEN, image=si("norton_2023_21st_century_art.png"))
    content_slide(prs, "BUILT ON QUIZZAP", "From Quiz App to Full SIS", [
        "Evolved from an auto-grading quiz app with student tracking…",
        "…into a fully-fledged Student Information System (SIS)",
        ("13 years", "supporting colleges & universities across the U.S., Canada, Hong Kong, and the West Indies"),
    ], GREEN, image=si("quizzap_gold_medal.png"))
    content_slide(prs, "THE HEALTH BRIDGE", "Your Data, Connected", [
        ("Interface to iOS Health", "Android features soon"),
        "Steps & energy, oxygen concentration, heart-rate variability, and more",
        "Configured to share all health data, enabling full interaction",
    ], GREEN, image=ei("lsb_runner_00001.gif"))
    content_slide(prs, "PARTNERSHIP", "Balanced Mind Wellness (BMW)", [
        "Analyze and interpret BMW reports using bio-energetics and more",
        "Help clients devise — and stick to — specific habits and routines",
        ("The path to lasting wellness", "one good habit at a time"),
    ], GREEN, image=ei("BMW-brain.png"))
    closing_slide(prs, "Mind. Body. Energy.", "lifesciencebalance.co", GREEN)
    save(prs, "SilverMedal_Healthcare_v2026.pptx")


def build_overview():
    prs = new_prs()
    ei = lambda n: img(IMG, "extra", n)
    ci = lambda n: img(IMG, "redrock", n)
    si = lambda n: img(IMG, "silvermedal", n)
    li = lambda n: img(LSB, n)
    title_slide(
        prs, "Silver Medal", "A CULTURE OF COMMUNITY & CREATIVITY",
        "silvermedal.net", BLUE, image=FOUR_ROBES)
    content_slide(prs, "OUR MISSION", "An Ecosystem of Tools", [
        "Tools that help people create, build, and improve their lives",
        "Free training and resources across our sites",
        ("A labor of love", "Dave is building this for his loved ones — to create jobs, opportunities, and to strengthen those in need"),
    ], BLUE, image=si("silvermedal_circular_medallion_0.png"))
    content_slide(prs, "THE ECOSYSTEM", "Three Worlds, One Roof", [
        ("Coding / Reasoning / Developing", "redrockcode.com"),
        ("Artwork / Healing / Wellness", "lifesciencebalance.co"),
        ("Writing / Publishing / Collaborating", "writersguild.online"),
    ], BLUE, image=si("silvermedal_four_robes_6.png"))
    content_slide(prs, "PILLAR ONE", "Red Rock Code", [
        "Project-based ES6+ app development — ports to iOS/Android as a PWA",
        "MVC & OOP from the ground up; the signature Workbench projects",
        ("CS 106A+ → Robotics", "Computer Science accreditation pending"),
    ], RED, image=ci("rrcc_shots_1.png"))
    content_slide(prs, "PILLAR TWO", "Life Science Balance", [
        "Artwork · Healing · Wellness — built on the Quizzap SIS",
        "Free three-part certification, NCLEX-RN 2026, new Graphic Arts course",
        ("Health Bridge", "+ Balanced Mind Wellness (BMW) partnership"),
    ], GREEN, image=ei("mandala_luminescent_1.png"))
    content_slide(prs, "PILLAR THREE", "Writers Guild", [
        "Collaborate, publish, sell — every author their own brand & domain",
        "Classics reborn (Morals & Dogma) + new works (The 144,000)",
        ("AI co-authorship", "and the Silver-Age Renaissance Series (SARS)"),
    ], PURPLE, image=ei("writers_guild_firefly_set5.png"))
    content_slide(prs, "ON THE HORIZON", "New Sites, Shared Kernel", [
        ("HeroBlue · heroblue.org", "Jurisprudence & Advocacy · Legal · 5 modules"),
        ("Coinbase Camp · coinbasecamp.org", "SIE, Series 7/63/65/66 · 5 modules + traderApp"),
        ("Quizzap · quizzap.org", "the kernel behind every course; licensing & open-source"),
    ], TEAL, image=si("quizzap_gold_medal.png"))
    content_slide(prs, "OUR PATH", "Building Toward Accreditation", [
        "Certifications for every completed Code Camp term (1000–3090)",
        "CS 106A + 108 materials and the Robotics extension in the works",
        ("Seeking accreditation", "in the State of Utah"),
    ], BLUE, image=si("silvermedal_circular_medallion_2.png"))
    closing_slide(prs, "Community · Creativity · Calendar", "silvermedal.net", BLUE)
    save(prs, "SilverMedal_Overview_v2026.pptx")


def save(prs, name):
    os.makedirs(OUT, exist_ok=True)
    path = os.path.join(OUT, name)
    prs.save(path)
    print(f"  wrote {os.path.relpath(path, ROOT)}  ({len(prs.slides._sldIdLst)} slides)")


def main():
    print("Silver Medal — building v2026 decks")
    build_coding()
    build_writing()
    build_healthcare()
    build_overview()
    print("Done.")


if __name__ == "__main__":
    main()
